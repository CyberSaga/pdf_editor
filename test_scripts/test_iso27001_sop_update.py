import subprocess
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SCRIPT_PATH = ROOT / "build" / "iso27001_sop_update" / "update_iso27001_sop.py"
SOURCE_DECK = next(ROOT.glob("*.pptx"))
OUTPUT_DECK = ROOT / "ISO27001審查項目SOP_附件加密與隨身碟加密.pptx"


def extract_texts(presentation):
    texts = []
    for slide in presentation.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip().replace("\n", " | "))
        texts.append(slide_texts)
    return texts


def iter_runs(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                yield run


def test_updated_iso27001_sop_deck_contains_new_encryption_section(tmp_path):
    generated_deck = tmp_path / OUTPUT_DECK.name

    subprocess.run(
        [
            sys.executable,
            str(SCRIPT_PATH),
            "--source",
            str(SOURCE_DECK),
            "--output",
            str(generated_deck),
        ],
        check=True,
        cwd=ROOT,
    )

    presentation = Presentation(str(generated_deck))
    slides = extract_texts(presentation)

    assert len(presentation.slides) == 80
    assert any("電子郵件附件加密／隨身碟加密" in " ".join(texts) for texts in slides)
    assert any(texts[0] == "電子郵件附件加密" for texts in slides if texts)
    assert any("for 系統 Windows 10 / Windows 11" in texts[0] for texts in slides if texts)
    assert any("BitLocker" in " ".join(texts) for texts in slides)
    assert any("Keka" in " ".join(texts) for texts in slides)

    for slide in list(presentation.slides)[59:]:
        content_box = next(
            (
                shape
                for shape in slide.shapes
                if getattr(shape, "name", "") == "Content Placeholder 2"
            ),
            None,
        )
        reminder_box = next(
            (
                shape
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip().startswith("作業提醒：")
            ),
            None,
        )

        if not content_box or not reminder_box:
            continue

        content_runs = list(iter_runs(content_box.text_frame))
        reminder_runs = list(iter_runs(reminder_box.text_frame))

        assert content_runs
        assert all(run.font.size and run.font.size.pt == 36 for run in content_runs)
        assert all(run.font.name in {"標楷體", "BiauKai"} for run in content_runs)
        assert not all(
            run.font.size and run.font.size.pt == 36 and run.font.name in {"標楷體", "BiauKai"}
            for run in reminder_runs
        )
