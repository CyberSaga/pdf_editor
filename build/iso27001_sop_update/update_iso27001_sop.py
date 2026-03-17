from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt


DEFAULT_OUTPUT_NAME = "ISO27001審查項目SOP_附件加密與隨身碟加密.pptx"


@dataclass(frozen=True)
class StepSlide:
    title: str
    detail: str
    image_hint: str = "示意圖（請貼操作畫面）"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=Path(DEFAULT_OUTPUT_NAME))
    return parser.parse_args()


def clear_placeholder(placeholder) -> None:
    if not hasattr(placeholder, "text_frame"):
        return
    placeholder.text = ""
    text_frame = placeholder.text_frame
    text_frame.clear()


def set_placeholder_text(placeholder, text: str, size: int | None = None) -> None:
    placeholder.text = text
    paragraph = placeholder.text_frame.paragraphs[0]
    if size is not None:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def add_section_slide(presentation: Presentation, title: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    set_placeholder_text(slide.shapes.title, title)
    if len(slide.placeholders) > 1:
        clear_placeholder(slide.placeholders[1])


def add_platform_slide(presentation: Presentation, title: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[2])
    set_placeholder_text(slide.shapes.title, title)
    if len(slide.placeholders) > 1:
        clear_placeholder(slide.placeholders[1])


def style_text_frame(text_frame, font_size: int, font_name: str | None = None) -> None:
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            if font_name:
                run.font.name = font_name


def add_step_slide(presentation: Presentation, step: StepSlide) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    set_placeholder_text(slide.shapes.title, step.title, size=32)

    body = slide.placeholders[1]
    body.left = 838200
    body.top = 1800878
    body.width = 4200000
    body.height = 1200000
    body.text = step.detail
    style_text_frame(body.text_frame, 36, "標楷體")

    image_left = 5480000
    image_top = 1500000
    image_width = 5300000
    image_height = 4300000
    image_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        image_left,
        image_top,
        image_width,
        image_height,
    )
    image_box.fill.solid()
    image_box.fill.fore_color.rgb = RGBColor(248, 249, 251)
    image_box.line.color.rgb = RGBColor(180, 187, 195)
    image_box.line.width = Pt(1.5)

    image_text = image_box.text_frame
    image_text.clear()
    paragraph = image_text.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = step.image_hint
    run.font.size = Pt(20)
    image_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    note_top = 3150000
    note_box = slide.shapes.add_textbox(838200, note_top, 4200000, 1700000)
    note_frame = note_box.text_frame
    note_frame.clear()
    note_paragraph = note_frame.paragraphs[0]
    note_paragraph.alignment = PP_ALIGN.LEFT
    note_run = note_paragraph.add_run()
    note_run.text = "作業提醒：密碼不得與附件同信寄送，請改以電話、簡訊或其他獨立管道通知收件人。"
    note_run.font.size = Pt(20)
    note_run.font.color.rgb = RGBColor(90, 90, 90)
    note_frame.word_wrap = True


def update_toc_slide(presentation: Presentation) -> None:
    toc_slide = presentation.slides[1]
    right_column = toc_slide.shapes[3]
    right_column.text = "電子郵件附件加密／隨身碟加密----------------58"
    style_text_frame(right_column.text_frame, 28)


def append_new_section(presentation: Presentation) -> None:
    add_section_slide(presentation, "電子郵件附件加密／隨身碟加密")

    add_section_slide(presentation, "電子郵件附件加密")
    add_platform_slide(presentation, "for 系統 Windows 10 / Windows 11")
    for step in [
        StepSlide(
            "1. 右鍵點選附件檔案後，選擇「7-Zip」→「加入壓縮檔」",
            "以下以公司核准之 7-Zip 建立加密壓縮檔為例；Win10 與 Win11 操作相同。",
        ),
        StepSlide(
            "2. 於壓縮格式選擇「zip」或公司規範格式",
            "若收件端未安裝 7-Zip，優先使用 `zip`；如公司另有規範，請依規範格式建立。",
        ),
        StepSlide(
            "3. 在 Encryption 區塊輸入密碼，並設定「AES-256」",
            "建議同時勾選「加密檔名（Encrypt file names）」；密碼需符合公司密碼政策，不得使用弱密碼。",
        ),
        StepSlide(
            "4. 確認已產生加密附件後，再將壓縮檔附加到電子郵件",
            "寄出前請自行測試可否正常開啟；密碼須以電話、簡訊或其他獨立通道通知收件人。",
        ),
    ]:
        add_step_slide(presentation, step)

    add_platform_slide(presentation, "for 系統 macOS")
    for step in [
        StepSlide(
            "1. 開啟 Keka，將要寄送的附件先整理到同一資料夾",
            "若尚未安裝 Keka，請先依公司核准流程安裝；Keka 可建立帶密碼的壓縮附件。",
        ),
        StepSlide(
            "2. 在 Keka 設定壓縮格式，並輸入附件密碼",
            "建議依收件端需求選擇 `zip` 或公司規範格式；密碼需符合公司政策，不得使用弱密碼。",
        ),
        StepSlide(
            "3. 將資料夾拖曳到 Keka 視窗，產生加密壓縮檔",
            "壓縮完成後請先自行測試是否可正常解壓縮；必要時可重新設定更清楚的檔名。",
        ),
        StepSlide(
            "4. 將加密壓縮檔附加到電子郵件，並以獨立管道通知密碼",
            "寄送前請確認附件為 Keka 產生的加密壓縮檔；密碼須另外透過電話、簡訊或其他獨立通道提供。",
        ),
    ]:
        add_step_slide(presentation, step)

    add_section_slide(presentation, "隨身碟加密")
    add_platform_slide(presentation, "for 系統 Windows 10 / Windows 11")
    for step in [
        StepSlide(
            "1. 插入隨身碟後，搜尋並開啟「管理 BitLocker」",
            "請確認選到的是可攜式磁碟機（Removable data drives），避免誤操作到系統磁碟。",
        ),
        StepSlide(
            "2. 在隨身碟旁點選「開啟 BitLocker」，並選擇以密碼解鎖",
            "請輸入符合公司密碼政策的密碼；若畫面提供智慧卡解鎖，非必要時可先不使用。",
        ),
        StepSlide(
            "3. 依畫面指示備份恢復金鑰，再開始加密",
            "恢復金鑰請依公司規範保存，不可與隨身碟一同存放；加密期間請勿直接拔除裝置。",
        ),
        StepSlide(
            "4. 等待加密完成後重新插入隨身碟，確認需輸入密碼才可存取",
            "完成驗證後再複製正式資料；若設備需交付他人，請以獨立管道傳遞解鎖密碼。",
        ),
    ]:
        add_step_slide(presentation, step)

    add_platform_slide(presentation, "for 系統 macOS")
    for step in [
        StepSlide(
            "1. 先備份隨身碟資料，再開啟「磁碟工具程式」",
            "macOS 對可攜式媒體啟用加密時，通常會重新格式化；請先確認資料已完整備份。",
        ),
        StepSlide(
            "2. 選擇「顯示所有裝置」，點選目標隨身碟後按「清除」",
            "請再次核對磁碟名稱與容量，避免誤清除其他外接磁碟或內建磁碟。",
        ),
        StepSlide(
            "3. 格式選擇加密格式，例如「APFS（加密）」或公司核准格式，並設定密碼",
            "若需相容舊版 macOS，可評估使用「Mac OS 擴充格式（加密）」；密碼需符合公司政策。",
        ),
        StepSlide(
            "4. 完成清除後重新掛載隨身碟，確認需輸入密碼才能開啟，再回存資料",
            "驗證完成後再放入正式資料；若需交付他人，密碼應以電話、簡訊或其他獨立管道提供。",
        ),
    ]:
        add_step_slide(presentation, step)


def main() -> None:
    args = parse_args()
    presentation = Presentation(str(args.source))
    update_toc_slide(presentation)
    append_new_section(presentation)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(args.output))


if __name__ == "__main__":
    main()
